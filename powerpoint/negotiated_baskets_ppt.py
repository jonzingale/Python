from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

p = Presentation()
p.slide_width = Inches(16)
p.slide_height = Inches(9)

# Title Slide
slide = p.slides.add_slide(p.slide_layouts[0])
slide.shapes.title.text = "Negotiated Baskets Trading Application"
slide.placeholders[1].text = "High-Level Specification"

# Application Overview
slide = p.slides.add_slide(p.slide_layouts[1])
slide.shapes.title.text = "Application Overview"
content = slide.placeholders[1]
content.text = (
    "Streamlines negotiating baskets of securities.\n"
    "Spreadsheet-like workflow for traders."
)

# System Architecture
slide = p.slides.add_slide(p.slide_layouts[1])
slide.shapes.title.text = "System Architecture"
content = slide.placeholders[1]
content.text = (
    "Technologies: Python (backend), Vue.js (frontend), IKIT and Markit (remote DB).\n"
    "External Systems: Bloomberg (market data), Markit (reference data), IKIT (logged data)."
)

# Inputs & Outputs Slide
slide = p.slides.add_slide(p.slide_layouts[1])
slide.shapes.title.text = "Inputs & Outputs"
content = slide.placeholders[1]
content.text = (
    "Inputs:\n"
    " Inventory: user may paste data or upload as CSV file\n"
    " Bloomberg: fetches unheld security data and for updating with market data (real-time or daily)\n"
    " Markit: provides held security reference details—daily updates are sufficient\n"
    "\n"
    "Outputs:\n"
    " Basket file: final ISINs and corresponding weights ready for execution"
)

# User Interface Concept
slide = p.slides.add_slide(p.slide_layouts[1])
slide.shapes.title.text = "User Interface Concept"
content = slide.placeholders[1]
content.text = (
    "Spreadsheet-like grid (Excel users).\n"
    "Filtering by attributes (sector, issuer, etc.).\n"
    "Highlighting flagged items.\n"
    "Hiding unused columns\n"
    "Parameterizable fields, locked reference data."
)

# Backend Modules slide
slide = p.slides.add_slide(p.slide_layouts[1])
slide.shapes.title.text = "Backend Modules Overview"
content = slide.placeholders[1]
content.text = (
    "Core backend modules for the application:\n"
    " Database:\n"
    "  - Security master: inventory, reference fund, model fund\n"
    "  - Logging: actions, data changes\n"
    " Inventory uploader:\n"
    "  - Supports pastable form fields and CSV upload\n"
    " Basket file downloader:\n"
    "  - Exports final basket (ISINs/weights) for execution\n"
    " Portfolio shape:\n"
    "  - User-editable model weights, reference weights, calculated differences\n"
    " Constraint tables:\n"
    "  - User sets filters and rules to include/exclude securities\n"
)

# Module Diagram
slide = p.slides.add_slide(p.slide_layouts[5])
slide.shapes.title.text = "Application Workflow Diagram"

# Element positions as {label, x, y}
boxes = [
    # Data sources
    {"label": "BBG\n(System Data)", "x": 0.5, "y": 1.0},
    {"label": "Markit\n(System Data)", "x": 0.5, "y": 2.2},
    {"label": "Inventory\n(User Input)", "x": 0.5, "y": 3.4},

    # Central data table
    {"label": "Security Master", "x": 3.2, "y": 2.2},

    # Controls & feedback
    {"label": "Portfolio Shape", "x": 6.2, "y": 1.2},
    {"label": "Constraints Table", "x": 6.2, "y": 3.2},

    # Basket & output
    {"label": "Basket View\n(securities for trade)", "x": 9.2, "y": 2.2},
    {"label": "Output\nCSV Basket File", "x": 12.0, "y": 2.2},
]

# Draw boxes
shape_map = {}
for b in boxes:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(b["x"]), Inches(b["y"]),
        Inches(2.2), Inches(1.1)
    )
    shape.text = b["label"]
    shape_map[b["label"]] = shape

# Connectors/arrows (start, end) by label
connectors = [
    # Data sources to Security Master
    ("BBG\n(System Data)", "Security Master"),
    ("Markit\n(System Data)", "Security Master"),
    ("Inventory\n(User Input)", "Security Master"),
    # Security Master to Portfolio Shape and Constraints Table
    ("Security Master", "Portfolio Shape"),
    ("Security Master", "Constraints Table"),
    # Feedback arrows (Portfolio Shape <-> Security Master)
    ("Portfolio Shape", "Security Master"),
    # Constraints to Portfolio Shape
    ("Constraints Table", "Portfolio Shape"),
    # Portfolio Shape to Basket View
    ("Portfolio Shape", "Basket View\n(securities for trade)"),
    # Constraints Table to Basket View
    ("Constraints Table", "Basket View\n(securities for trade)"),
    # Basket View to Output
    ("Basket View\n(securities for trade)", "Output\nCSV Basket File"),
]

# Draw arrows
for start, end in connectors:
    start_shape = shape_map[start]
    end_shape = shape_map[end]
    start_x = (start_shape.left + start_shape.width)
    start_y = (start_shape.top + start_shape.height / 2)
    end_x = end_shape.left
    end_y = (end_shape.top + end_shape.height / 2)
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        start_x, start_y,
        end_x, end_y
    )

# Optional: adjust arrow directions, styles, box colors, and add explanation

# Next Steps / Open Questions
slide = p.slides.add_slide(p.slide_layouts[1])
slide.shapes.title.text = "Next Steps / Open Questions"
content = slide.placeholders[1]
content.text = (
    "Review with stakeholders.\n"
    " Does the workflow make sense?\n"
    " Is this approach scalable and aligned with institutional goals?\n"
    "Confirm UI requirements, including:\n"
    " Research spreadsheet-like workflow options: Jspreadsheet, DHTMLX, Webix\n"
    " Schedule UI design meeting\n"
)

p.save("negotiated_baskets_spec.pptx")
